"""
投資ニューススライド自動生成システム
バフェットかおるチャンネル用 - Gemini API + python-pptx

毎朝8時に実行し、その日のバズりそうな投資・経済ニュースを
PowerPointスライドと読み上げ用カンペとして自動生成する。
"""

import json
import os
import sys
from datetime import datetime
from pathlib import Path

from dotenv import load_dotenv
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# ── 設定 ──────────────────────────────────────────────
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

TODAY = datetime.now().strftime("%Y年%m月%d日")
TODAY_FILE = datetime.now().strftime("%Y%m%d")


# ── 1. Gemini API でニュース取得 ──────────────────────
def fetch_news_topics() -> list[dict]:
    """Gemini APIを使って、本日の注目投資ニュース3トピックを取得する。"""
    if not GEMINI_API_KEY:
        print("エラー: GEMINI_API_KEY が設定されていません。.env ファイルを確認してください。")
        sys.exit(1)
    client = genai.Client(api_key=GEMINI_API_KEY)

    prompt = f"""あなたは投資・経済ニュースの専門アナリストです。
今日は{TODAY}です。

以下の条件で、今日バズりそうな投資・経済ニュースを**3つ**選んでください。

【選定基準】
- 米国市場（NYSE/NASDAQ）と日本市場（東証）の両方からバランスよく選ぶ
- 個人投資家の関心が高いトピック（株価の大きな動き、決算、政策変更、為替など）
- YouTubeの投資系チャンネルで視聴回数が伸びそうなテーマ

【各トピックに含める情報】
1. title: トピックのタイトル（YouTube動画のスライドに使える簡潔なもの、20文字以内）
2. market: "米国市場" または "日本市場"
3. summary: 概要（何が起きたのか、背景を含めて3〜4文で説明）
4. impact: 投資家への影響（個人投資家が知るべきポイントを2〜3文で説明）
5. action_plan: 今日のアクションプラン（具体的に何をすべきか、2〜3文で説明）

【出力形式】
以下のJSON形式で出力してください。JSON以外のテキストは含めないでください。
```json
[
  {{
    "title": "...",
    "market": "...",
    "summary": "...",
    "impact": "...",
    "action_plan": "..."
  }},
  ...
]
```
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
    )

    text = response.text.strip()
    # コードブロックが含まれる場合は除去
    if text.startswith("```"):
        text = text.split("\n", 1)[1]
    if text.endswith("```"):
        text = text.rsplit("```", 1)[0]
    text = text.strip()

    topics = json.loads(text)
    if not isinstance(topics, list) or len(topics) == 0:
        raise ValueError("Gemini APIから有効なトピックリストを取得できませんでした。")
    return topics


# ── 2. スライド作成 ────────────────────────────────────

# カラーパレット
COLOR_BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0xD4, 0xAA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
COLOR_MARKET_US = RGBColor(0x44, 0x88, 0xFF)
COLOR_MARKET_JP = RGBColor(0xFF, 0x44, 0x44)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color: RGBColor):
    """スライド背景色を設定する。"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Meiryo"):
    """テキストボックスを追加するヘルパー。"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_shape_rect(slide, left, top, width, height, fill_color: RGBColor):
    """色付き矩形を追加する。"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def create_title_slide(prs: Presentation):
    """タイトルスライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, COLOR_BG_DARK)

    # アクセントライン上部
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)

    # チャンネル名
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                "バフェットかおるの投資ニュース",
                font_size=44, color=COLOR_ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 日付
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                f"{TODAY} 朝の市場チェック",
                font_size=28, color=COLOR_WHITE, bold=False,
                alignment=PP_ALIGN.CENTER)

    # サブタイトル
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                "今日バズる投資ニュース TOP 3",
                font_size=24, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # アクセントライン下部
    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)


def create_topic_slide(prs: Presentation, topic: dict, index: int):
    """各トピックのスライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG_DARK)

    market_color = COLOR_MARKET_US if "米国" in topic["market"] else COLOR_MARKET_JP

    # ── ヘッダー部分 ──
    # トピック番号バッジ
    add_shape_rect(slide, Inches(0.5), Inches(0.4), Inches(0.9), Inches(0.7), COLOR_ACCENT)
    add_textbox(slide, Inches(0.5), Inches(0.42), Inches(0.9), Inches(0.7),
                f"#{index}", font_size=32, color=COLOR_BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)

    # マーケットバッジ
    add_shape_rect(slide, Inches(1.6), Inches(0.4), Inches(2.0), Inches(0.7), market_color)
    add_textbox(slide, Inches(1.6), Inches(0.42), Inches(2.0), Inches(0.7),
                topic["market"], font_size=20, color=COLOR_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # タイトル
    add_textbox(slide, Inches(4.0), Inches(0.35), Inches(8.5), Inches(0.8),
                topic["title"], font_size=32, color=COLOR_WHITE, bold=True)

    # 区切り線
    add_shape_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.03), COLOR_ACCENT)

    # ── 3カラムレイアウト ──
    sections = [
        ("概要", topic["summary"], Inches(0.5)),
        ("投資家への影響", topic["impact"], Inches(4.7)),
        ("今日のアクションプラン", topic["action_plan"], Inches(8.9)),
    ]

    for label, content, left in sections:
        # セクションラベル
        add_shape_rect(slide, left, Inches(1.7), Inches(3.8), Inches(0.55), COLOR_ACCENT)
        add_textbox(slide, left, Inches(1.72), Inches(3.8), Inches(0.55),
                    f"  {label}", font_size=18, color=COLOR_BG_DARK, bold=True)

        # セクション内容
        add_textbox(slide, left, Inches(2.5), Inches(3.8), Inches(4.5),
                    content, font_size=16, color=COLOR_LIGHT_GRAY)

    # フッターライン
    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)


# ── 投資分析スライド ────────────────────────────────────

COLOR_GOLD = RGBColor(0xFF, 0xD7, 0x00)
COLOR_SECTION_BG = RGBColor(0x22, 0x22, 0x3A)

INVESTMENT_STOCKS = [
    {
        "name": "日東富士製粉 (2003)",
        "stability": "2026年までの中計で「累進配当」を導入。自己資本比率も高く、下限配当を設定できるほどキャッシュが潤沢。",
        "reason": "パン、麺類、お菓子などの原料「小麦粉」の国内大手。食卓に欠かせない究極のディフェンシブ業種。",
        "keyword": "食べる（小麦粉）",
    },
    {
        "name": "ブリヂストン (5108)",
        "stability": "世界首位級のタイヤメーカー。2026年度も増収増益を見込み、1,500億円規模の自社株買いを行うほど財務力がある。",
        "reason": "「移動」を支えるインフラ企業。EVになってもタイヤは必ず消耗・交換が必要で、社会の物流と移動に不可欠。",
        "keyword": "移動する（タイヤ）",
    },
    {
        "name": "積水ハウス (1928)",
        "stability": "国内外で強固な顧客基盤を持ち、配当性向40%以上を維持。米国市場の成長も取り込み安定した収益構造。",
        "reason": "「衣・食・住」の「住」のトップランナー。良質な住宅供給は家族の安全と幸せな暮らしを支える社会の土台。",
        "keyword": "住まう（住宅）",
    },
]


def create_investment_title_slide(prs: Presentation):
    """投資分析セクションのタイトルスライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG_DARK)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), COLOR_GOLD)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                f"{TODAY}、私が「生活の土台」に投資した理由",
                font_size=38, color=COLOR_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                "本日、3つの日本株に投資しました",
                font_size=26, color=COLOR_WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                "日東富士製粉 ・ ブリヂストン ・ 積水ハウス",
                font_size=24, color=COLOR_ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(5.4), Inches(9), Inches(1.2),
                "「どんな時代になっても、人が生きていく上で絶対に欠かせないもの」を提供している企業に投資する",
                font_size=18, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_GOLD)


def create_stock_slide(prs: Presentation, stock: dict, index: int):
    """個別銘柄の分析スライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG_DARK)

    # ── ヘッダー ──
    add_shape_rect(slide, Inches(0.5), Inches(0.4), Inches(0.9), Inches(0.7), COLOR_GOLD)
    add_textbox(slide, Inches(0.5), Inches(0.42), Inches(0.9), Inches(0.7),
                f"#{index}", font_size=32, color=COLOR_BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.8), Inches(0.35), Inches(10), Inches(0.8),
                stock["name"], font_size=32, color=COLOR_WHITE, bold=True)

    # キーワードバッジ
    add_shape_rect(slide, Inches(1.8), Inches(1.1), Inches(3.5), Inches(0.5), COLOR_ACCENT)
    add_textbox(slide, Inches(1.8), Inches(1.12), Inches(3.5), Inches(0.5),
                f"  {stock['keyword']}", font_size=16, color=COLOR_BG_DARK, bold=True)

    # 区切り線
    add_shape_rect(slide, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.03), COLOR_GOLD)

    # ── 2カラムレイアウト ──
    # 財務安定性
    add_shape_rect(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(0.55), COLOR_GOLD)
    add_textbox(slide, Inches(0.5), Inches(2.32), Inches(5.8), Inches(0.55),
                "  財務安定性", font_size=18, color=COLOR_BG_DARK, bold=True)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.8), Inches(3.5),
                stock["stability"], font_size=16, color=COLOR_LIGHT_GRAY)

    # 生活に必要な理由
    add_shape_rect(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(0.55), COLOR_ACCENT)
    add_textbox(slide, Inches(7.0), Inches(2.32), Inches(5.8), Inches(0.55),
                "  生活に必要な理由", font_size=18, color=COLOR_BG_DARK, bold=True)
    add_textbox(slide, Inches(7.0), Inches(3.1), Inches(5.8), Inches(3.5),
                stock["reason"], font_size=16, color=COLOR_LIGHT_GRAY)

    # フッターライン
    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_GOLD)


def create_inflation_slide(prs: Presentation):
    """インフレ対策スライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG_DARK)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), COLOR_GOLD)

    add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                "インフレ負けしない資産を持つ意味",
                font_size=34, color=COLOR_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.03), COLOR_GOLD)

    add_textbox(slide, Inches(1), Inches(1.9), Inches(11), Inches(1.5),
                "インフレとは「モノの価値が上がり、お金の価値が下がること」です。\n"
                "銀行に現金を置いているだけでは、購買力が日々削られていくリスクがあります。",
                font_size=18, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.LEFT)

    add_shape_rect(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(2.2), COLOR_SECTION_BG)
    add_textbox(slide, Inches(1.3), Inches(3.9), Inches(10.7), Inches(2.0),
                "「生活に不可欠なサービスを提供する企業の株」を持つことは、\n"
                "物価上昇に合わせて製品価格を適正化できる「実物資産」を持つことを意味します。\n\n"
                "企業が稼ぎ続ける限り、配当や株価上昇という形で、\n"
                "私たちの資産をインフレの波から守ってくれるのです。",
                font_size=18, color=COLOR_WHITE,
                alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
                "お金を「貯める」から、価値を生む場所に「置く」へ。",
                font_size=24, color=COLOR_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_GOLD)


def create_ending_slide(prs: Presentation):
    """エンディングスライドを作成する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG_DARK)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
                "チャンネル登録・高評価お願いします！",
                font_size=36, color=COLOR_ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                "コメント欄であなたの投資戦略を教えてください",
                font_size=24, color=COLOR_WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
                "※投資は自己責任でお願いします。本動画は情報提供を目的としており、特定の銘柄の売買を推奨するものではありません。",
                font_size=14, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)


def generate_slides(topics: list[dict]) -> Path:
    """全スライドを生成し、.pptxファイルとして保存する。"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_title_slide(prs)
    for i, topic in enumerate(topics, 1):
        create_topic_slide(prs, topic, i)

    # 投資分析スライド
    create_investment_title_slide(prs)
    for i, stock in enumerate(INVESTMENT_STOCKS, 1):
        create_stock_slide(prs, stock, i)
    create_inflation_slide(prs)

    create_ending_slide(prs)

    output_path = OUTPUT_DIR / f"investment_news_{TODAY_FILE}.pptx"
    prs.save(str(output_path))
    return output_path


# ── 3. 読み上げ用カンペ生成 ─────────────────────────────

def generate_script(topics: list[dict]) -> Path:
    """読み上げ用カンペをテキストファイルとして生成する。"""
    lines = []
    lines.append("=" * 60)
    lines.append(f"  バフェットかおるの投資ニュース - {TODAY}")
    lines.append(f"  読み上げ用カンペ（台本）")
    lines.append("=" * 60)
    lines.append("")

    # オープニング
    lines.append("【オープニング】")
    lines.append("")
    lines.append("はい、皆さんおはようございます！バフェットかおるです。")
    lines.append(f"今日は{TODAY}、朝の投資ニュースチェックをやっていきましょう。")
    lines.append("今日も注目のニュースを3つピックアップしました。")
    lines.append("それでは早速いきましょう！")
    lines.append("")
    lines.append("-" * 60)

    for i, topic in enumerate(topics, 1):
        lines.append("")
        lines.append(f"【トピック{i}】{topic['title']}（{topic['market']}）")
        lines.append("")
        lines.append(f"  ▼ 概要")
        lines.append(f"  {topic['summary']}")
        lines.append("")
        lines.append(f"  ▼ 投資家への影響")
        lines.append(f"  {topic['impact']}")
        lines.append("")
        lines.append(f"  ▼ 今日のアクションプラン")
        lines.append(f"  {topic['action_plan']}")
        lines.append("")

        if i < len(topics):
            lines.append("  → 「続いてはこちらのニュースです。」")
        lines.append("-" * 60)

    # 投資分析セクション
    lines.append("")
    lines.append("=" * 60)
    lines.append(f"  {TODAY}、私が「生活の土台」に投資した理由")
    lines.append("=" * 60)
    lines.append("")
    lines.append("  → 「さて、ここからは本日私が実際に投資した3銘柄についてお話しします。」")
    lines.append("")
    lines.append("本日、3つの日本株に投資しました。")
    lines.append("日東富士製粉、ブリヂストン、そして積水ハウスです。")
    lines.append("")
    lines.append("なぜこの3銘柄なのか？ 理由はシンプルです。")
    lines.append("「どんな時代になっても、人が生きていく上で絶対に欠かせないもの」を提供しているからです。")
    lines.append("")
    lines.append("-" * 60)

    for i, stock in enumerate(INVESTMENT_STOCKS, 1):
        lines.append("")
        lines.append(f"【銘柄{i}】{stock['name']}　── {stock['keyword']}")
        lines.append("")
        lines.append(f"  ▼ 財務安定性")
        lines.append(f"  {stock['stability']}")
        lines.append("")
        lines.append(f"  ▼ 生活に必要な理由")
        lines.append(f"  {stock['reason']}")
        lines.append("")
        lines.append("-" * 60)

    lines.append("")
    lines.append("【インフレ対策としての株式投資】")
    lines.append("")
    lines.append("インフレとは「モノの価値が上がり、お金の価値が下がること」です。")
    lines.append("銀行に現金を置いているだけでは、購買力が日々削られていくリスクがあります。")
    lines.append("")
    lines.append("一方で、「生活に不可欠なサービスを提供する企業の株」を持つことは、")
    lines.append("物価上昇に合わせて製品価格を適正化できる「実物資産」を持つことを意味します。")
    lines.append("企業が稼ぎ続ける限り、配当や株価上昇という形で、私たちの資産をインフレの波から守ってくれるのです。")
    lines.append("")
    lines.append("お金を「貯める」から、価値を生む場所に「置く」へ。")
    lines.append("皆さんも自分の生活を支えてくれる身近な企業に投資してみてはいかがでしょうか。")
    lines.append("")
    lines.append("-" * 60)

    # エンディング
    lines.append("")
    lines.append("【エンディング】")
    lines.append("")
    lines.append("はい、以上が今日の注目ニュース3選でした。")
    lines.append("皆さんはどのニュースが一番気になりましたか？")
    lines.append("コメント欄であなたの投資戦略を教えてください。")
    lines.append("")
    lines.append("それでは、チャンネル登録と高評価、よろしくお願いします。")
    lines.append("バフェットかおるでした。また明日！")
    lines.append("")
    lines.append("※投資は自己責任でお願いします。")
    lines.append("  本動画は情報提供を目的としており、特定の銘柄の売買を推奨するものではありません。")
    lines.append("")
    lines.append("=" * 60)

    output_path = OUTPUT_DIR / f"script_{TODAY_FILE}.txt"
    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path


# ── メイン処理 ─────────────────────────────────────────

def main():
    print(f"[{datetime.now():%H:%M:%S}] 投資ニューススライド生成を開始します...")

    print(f"[{datetime.now():%H:%M:%S}] Gemini APIからニュースを取得中...")
    topics = fetch_news_topics()
    print(f"[{datetime.now():%H:%M:%S}] {len(topics)}件のトピックを取得しました。")

    for i, t in enumerate(topics, 1):
        print(f"  #{i} [{t['market']}] {t['title']}")

    print(f"[{datetime.now():%H:%M:%S}] スライドを生成中...")
    pptx_path = generate_slides(topics)
    print(f"[{datetime.now():%H:%M:%S}] スライド保存: {pptx_path}")

    print(f"[{datetime.now():%H:%M:%S}] 読み上げカンペを生成中...")
    script_path = generate_script(topics)
    print(f"[{datetime.now():%H:%M:%S}] カンペ保存: {script_path}")

    print(f"[{datetime.now():%H:%M:%S}] 完了！")


if __name__ == "__main__":
    main()
